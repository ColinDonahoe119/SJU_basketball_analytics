"""
Build transfers_analysis.pptx from the transfer portal category analysis findings.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

BASE = r"C:\Users\colin\OneDrive - Saint Joseph's University\Documents\SJU_Basketball_Analytics"

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
SJU_CRIMSON  = RGBColor(0x99, 0x00, 0x00)   # SJU red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x22, 0x22, 0x22)
LIGHT_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
MID_GRAY     = RGBColor(0xCC, 0xCC, 0xCC)
BLUE         = RGBColor(0x21, 0x96, 0xF3)
GREEN        = RGBColor(0x4C, 0xAF, 0x50)
ORANGE       = RGBColor(0xFF, 0x57, 0x22)
GOLD         = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def bg(slide, color):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_image(slide, path, l, t, w, h=None):
    if not os.path.exists(path):
        return
    if h:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))

def accent_bar(slide, color, height=0.08):
    """Thin colored bar across the top."""
    add_rect(slide, 0, 0, 13.33, height, fill_color=color)

def header_band(slide, color, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill_color=color)
    add_text(slide, title, 0.35, 0.12, 12.5, 0.7,
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.78, 12.5, 0.42,
                 font_size=16, bold=False, color=RGBColor(0xFF,0xDD,0xDD),
                 align=PP_ALIGN.LEFT)

def stat_card(slide, label, value, l, t, w=1.9, h=1.05,
              bg_color=SJU_CRIMSON, val_color=WHITE, lbl_color=RGBColor(0xFF,0xDD,0xDD)):
    add_rect(slide, l, t, w, h, fill_color=bg_color)
    add_text(slide, value, l+0.08, t+0.05, w-0.16, 0.6,
             font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, l+0.08, t+0.62, w-0.16, 0.38,
             font_size=11, bold=False, color=lbl_color, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h, title=None,
               bg_color=LIGHT_GRAY, title_color=DARK_GRAY, text_color=DARK_GRAY,
               font_size=13):
    add_rect(slide, l, t, w, h, fill_color=bg_color)
    yo = t + 0.12
    if title:
        add_text(slide, title, l+0.15, yo, w-0.3, 0.35,
                 font_size=13, bold=True, color=title_color)
        yo += 0.35
    for item in items:
        add_text(slide, f"  {item}", l+0.1, yo, w-0.2, 0.35,
                 font_size=font_size, color=text_color)
        yo += 0.34

def divider(slide, t, color=MID_GRAY):
    add_rect(slide, 0.35, t, 12.63, 0.03, fill_color=color)


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, DARK_GRAY)
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_GRAY)

# crimson gradient band
add_rect(slide, 0, 2.5, 13.33, 2.7, fill_color=SJU_CRIMSON)

add_text(slide, "SJU BASKETBALL ANALYTICS",
         0.5, 1.3, 12.33, 0.6,
         font_size=16, bold=True, color=RGBColor(0xBB,0xBB,0xBB),
         align=PP_ALIGN.CENTER)

add_text(slide, "Transfer Portal",
         0.5, 2.6, 12.33, 1.0,
         font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, "Player Category Analysis",
         0.5, 3.5, 12.33, 0.75,
         font_size=36, bold=False, color=RGBColor(0xFF,0xDD,0xDD),
         align=PP_ALIGN.CENTER)

add_text(slide, "Elite Rebounders  |  Elite Playmakers  |  Two-Way Engines",
         0.5, 4.45, 12.33, 0.5,
         font_size=15, bold=False, color=RGBColor(0xAA,0xAA,0xAA),
         align=PP_ALIGN.CENTER)

add_text(slide, "n = 3,840 transfer seasons  |  Metrics: barthag, WAB, BPM, oreb_rate, ast/to",
         0.5, 6.7, 12.33, 0.4,
         font_size=11, italic=True, color=RGBColor(0x88,0x88,0x88),
         align=PP_ALIGN.CENTER)


# ===========================================================================
# SLIDE 2 — Methodology / Thresholds
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Methodology", "How we segmented 3,840 transfer seasons into 3 categories")

# threshold cards
stat_card(slide, "oreb_rate threshold (75th pct)", "7.70", 0.5, 1.55, w=3.0, h=1.1,
          bg_color=BLUE)
stat_card(slide, "ast/to threshold (75th pct)", "1.45", 4.0, 1.55, w=3.0, h=1.1,
          bg_color=GREEN)
stat_card(slide, "Min MPG filter", "5.0", 7.5, 1.55, w=2.5, h=1.1,
          bg_color=DARK_GRAY)
stat_card(slide, "Min games filter", "5 games", 10.3, 1.55, w=2.5, h=1.1,
          bg_color=DARK_GRAY)

# category definitions
cats = [
    ("A  —  Elite Rebounder", BLUE,   "oreb_rate >= 7.70   AND   ast/to < 1.45",
     "924 players  |  Mostly bigs (C, PF/C, Wing F)"),
    ("B  —  Elite Playmaker", GREEN,  "ast/to >= 1.45   AND   oreb_rate < 7.70",
     "914 players  |  Guards & wings (Combo G, Scoring PG, Wing G)"),
    ("C  —  Two-Way Engine",  ORANGE, "oreb_rate >= 7.70   AND   ast/to >= 1.45",
     "46 players   |  Versatile forwards (Wing F, PF/C, C)"),
]
for i, (label, color, crit, note) in enumerate(cats):
    y = 2.9 + i * 1.35
    add_rect(slide, 0.5, y, 12.33, 1.18, fill_color=color)
    add_text(slide, label, 0.7, y+0.08, 4.5, 0.48, font_size=18, bold=True, color=WHITE)
    add_text(slide, crit,  5.4, y+0.08, 7.0, 0.48, font_size=15, bold=False, color=WHITE)
    add_text(slide, note,  0.7, y+0.62, 11.8, 0.38, font_size=12, italic=True,
             color=RGBColor(0xFF,0xEE,0xEE))


# ===========================================================================
# SLIDE 3 — Category A: Elite Rebounders
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, BLUE, "Category A  —  Elite Rebounders", "n = 924  |  oreb_rate >= 7.70  |  ast/to < 1.45")

# stat cards
stat_card(slide, "Avg oreb_rate",  "10.7",  0.35, 1.55, bg_color=BLUE)
stat_card(slide, "Avg ast/to",     "0.66",  2.45, 1.55, bg_color=BLUE)
stat_card(slide, "Avg BPM",        "+0.13", 4.55, 1.55, bg_color=BLUE)
stat_card(slide, "Avg ORTG",       "106.6", 6.65, 1.55, bg_color=BLUE)
stat_card(slide, "Avg MPG",        "18.7",  8.75, 1.55, bg_color=BLUE)
stat_card(slide, "Avg TS%",        "55.8",  10.85,1.55, bg_color=BLUE)

# position breakdown box
bullet_box(slide,
    ["C             —  413 (45%)",
     "PF/C          —  260 (28%)",
     "Wing F        —  136 (15%)",
     "Stretch 4     —   80  (9%)",
     "Wing G        —   35  (4%)"],
    0.35, 2.9, 3.8, 2.35,
    title="Position Breakdown", bg_color=LIGHT_GRAY, title_color=BLUE)

# key findings box
bullet_box(slide,
    ["Land at teams with avg barthag 0.548",
     "Barthag change vs previous school: -0.005",
     "(Tend to trade down for more minutes)",
     "Destination WAB avg: -6.7",
     "Stacking 2+: barthag rises to 0.575 (p=0.019)"],
    4.35, 2.9, 5.3, 2.35,
    title="Winning Impact", bg_color=LIGHT_GRAY, title_color=BLUE)

# top feeder schools
bullet_box(slide,
    ["Louisville — 10", "George Mason — 9", "Xavier — 9",
     "Vanderbilt — 8", "East Carolina — 8"],
    9.85, 2.9, 3.1, 2.35,
    title="Top Feeder Schools", bg_color=LIGHT_GRAY, title_color=BLUE)

add_text(slide,
    "Key Insight: Elite Rebounders are predominantly bigs who move laterally or slightly down "
    "in program strength. They provide a floor-spacing complement but show the smallest "
    "individual winning uplift of the three categories.",
    0.35, 5.45, 12.63, 0.9,
    font_size=13, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 4 — Category B: Elite Playmakers
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, GREEN, "Category B  —  Elite Playmakers", "n = 914  |  ast/to >= 1.45  |  oreb_rate < 7.70")

stat_card(slide, "Avg oreb_rate",  "2.6",   0.35, 1.55, bg_color=GREEN)
stat_card(slide, "Avg ast/to",     "2.01",  2.45, 1.55, bg_color=GREEN)
stat_card(slide, "Avg BPM",        "+0.55", 4.55, 1.55, bg_color=GREEN)
stat_card(slide, "Avg ORTG",       "105.7", 6.65, 1.55, bg_color=GREEN)
stat_card(slide, "Avg MPG",        "24.6",  8.75, 1.55, bg_color=GREEN)
stat_card(slide, "Avg TS%",        "51.8",  10.85,1.55, bg_color=GREEN)

bullet_box(slide,
    ["Combo G      —  288 (32%)",
     "Scoring PG   —  272 (30%)",
     "Wing G       —  188 (21%)",
     "Pure PG      —   85 (9%)",
     "Stretch 4    —   34 (4%)"],
    0.35, 2.9, 3.8, 2.35,
    title="Position Breakdown", bg_color=LIGHT_GRAY, title_color=GREEN)

bullet_box(slide,
    ["Land at teams with avg barthag 0.593",
     "Barthag change vs previous school: +0.064",
     "(Playmakers move UP in program quality)",
     "Destination WAB avg: -5.4",
     "Stacking 2+: barthag jumps to 0.659 (p<0.0001)"],
    4.35, 2.9, 5.3, 2.35,
    title="Winning Impact", bg_color=LIGHT_GRAY, title_color=GREEN)

bullet_box(slide,
    ["West Virginia — 9", "UT Rio Grande Valley — 8",
     "Arizona St. — 8", "Indiana St. — 8", "Rutgers — 8"],
    9.85, 2.9, 3.1, 2.35,
    title="Top Feeder Schools", bg_color=LIGHT_GRAY, title_color=GREEN)

add_text(slide,
    "Key Insight: Elite Playmakers log the most minutes (24.6 MPG) and are the most "
    "transferable asset in the portal. Programs that stack 2+ see a dramatic barthag jump "
    "from 0.544 to 0.659 — the largest stacking effect of any category.",
    0.35, 5.45, 12.63, 0.9,
    font_size=13, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 5 — Category C: Two-Way Engines
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, ORANGE, "Category C  —  Two-Way Engines", "n = 46  |  oreb_rate >= 7.70  AND  ast/to >= 1.45")

stat_card(slide, "Avg oreb_rate",  "9.97",  0.35, 1.55, bg_color=ORANGE)
stat_card(slide, "Avg ast/to",     "2.02",  2.45, 1.55, bg_color=ORANGE)
stat_card(slide, "Avg BPM",        "+2.55", 4.55, 1.55, bg_color=ORANGE)
stat_card(slide, "Avg ORTG",       "115.0", 6.65, 1.55, bg_color=ORANGE)
stat_card(slide, "Avg MPG",        "18.8",  8.75, 1.55, bg_color=ORANGE)
stat_card(slide, "Avg TS%",        "53.9",  10.85,1.55, bg_color=ORANGE)

bullet_box(slide,
    ["Wing F      —  13 (28%)",
     "PF/C        —  11 (24%)",
     "C           —   7 (15%)",
     "Wing G      —   7 (15%)",
     "Stretch 4   —   7 (15%)"],
    0.35, 2.9, 3.8, 2.35,
    title="Position Breakdown", bg_color=LIGHT_GRAY, title_color=ORANGE)

bullet_box(slide,
    ["Highest dest. barthag: 0.613 (rank #1)",
     "Barthag change vs previous school: +0.104",
     "(Biggest upward move of all categories)",
     "Destination WAB avg: -4.6",
     "Only 46 players — extremely rare profile"],
    4.35, 2.9, 5.3, 2.35,
    title="Winning Impact", bg_color=LIGHT_GRAY, title_color=ORANGE)

bullet_box(slide,
    ["La Salle — 2", "Southern Utah — 2",
     "Central Michigan — 1",
     "Georgia Southern — 1", "Davidson — 1"],
    9.85, 2.9, 3.1, 2.35,
    title="Top Feeder Schools", bg_color=LIGHT_GRAY, title_color=ORANGE)

add_text(slide,
    "Key Insight: Two-Way Engines are the rarest and most impactful players in the portal. "
    "BPM of +2.55 is 20x higher than Elite Rebounders. They are disproportionately versatile "
    "forwards, not guards, and land at the strongest programs of any category.",
    0.35, 5.45, 12.63, 0.9,
    font_size=13, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 6 — Winning Impact Comparison (charts)
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Winning Impact", "Destination team strength by category vs. baseline (other transfers)")

# Embed the two charts side by side
add_image(slide, f"{BASE}/chart_boxplot_team_success.png",  0.35, 1.4, 6.4)
add_image(slide, f"{BASE}/chart_barthag_change.png",        7.0,  1.4, 6.0)

add_text(slide,
    "All three categories land at significantly stronger programs than baseline transfers (p < 0.01).  "
    "Two-Way Engines show the largest upward move (+0.104 barthag); Elite Rebounders move slightly DOWN (-0.005).",
    0.35, 6.55, 12.63, 0.7,
    font_size=12, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 7 — Stacking Effect
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "The Stacking Effect", "Does acquiring 2+ players of the same category amplify team success?")

# data table visual
rows_data = [
    ("Category",          "1 Player\nbarthag", "2+ Players\nbarthag", "Difference", "p-value", "Verdict"),
    ("A - Elite Rebounder","0.518",             "0.575",               "+0.057",     "0.019",   "Significant"),
    ("B - Elite Playmaker","0.544",             "0.659",               "+0.115",     "<0.0001", "Highly Significant"),
    ("C - Two-Way Engine", "N/A",               "N/A",                 "N/A",        "N/A",     "Too few cases"),
]
col_widths = [3.0, 1.8, 1.9, 1.6, 1.4, 2.3]
col_x = [0.35]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w)

header_bg = SJU_CRIMSON
row_bgs   = [LIGHT_GRAY, WHITE, LIGHT_GRAY]

for ri, row in enumerate(rows_data):
    y = 1.55 + ri * 0.88
    h = 0.85
    for ci, (cell, cw, cx) in enumerate(zip(row, col_widths, col_x)):
        bg_c = header_bg if ri == 0 else row_bgs[(ri-1) % 2]
        add_rect(slide, cx, y, cw-0.04, h, fill_color=bg_c)
        fc = WHITE if ri == 0 else DARK_GRAY
        if ri > 0 and ci == 3 and cell.startswith("+"):
            fc = GREEN
        if ri > 0 and ci == 5 and "Highly" in cell:
            fc = GREEN
        add_text(slide, cell, cx+0.08, y+0.08, cw-0.2, h-0.12,
                 font_size=12 if ri > 0 else 11,
                 bold=(ri == 0), color=fc, align=PP_ALIGN.CENTER)

add_text(slide,
    "Takeaway: Stacking Elite Playmakers is the single highest-leverage portal strategy. "
    "Teams that acquired 2+ playmakers averaged a 0.659 barthag — deep tournament-contender territory.",
    0.35, 5.45, 12.63, 0.9,
    font_size=14, italic=True, color=DARK_GRAY)

add_image(slide, f"{BASE}/chart_category_ranking.png", 8.5, 1.55, 4.6)


# ===========================================================================
# SLIDE 8 — Category Scatter (visual)
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Player Segmentation Map", "oreb_rate vs ast/to — all 3,840 transfer seasons")

add_image(slide, f"{BASE}/chart_scatter_categories.png", 1.5, 1.35, 10.33)

add_text(slide,
    "Dashed lines mark the 75th-percentile thresholds. Two-Way Engines (orange, upper-right) "
    "are rare but represent the highest-impact transfer profile.",
    0.35, 6.7, 12.63, 0.6,
    font_size=12, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 9 — Position Breakdown (visual)
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Position Breakdown by Category", "Who is transferring — and where are they playing?")

add_image(slide, f"{BASE}/chart_position_breakdown.png", 0.6, 1.35, 12.13)

add_text(slide,
    "Elite Rebounders are dominated by traditional bigs. Elite Playmakers are almost entirely guards. "
    "Two-Way Engines are uniquely versatile — split across Wing F, PF/C, C, and Wing G.",
    0.35, 6.75, 12.63, 0.55,
    font_size=12, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 10 — Summary & Rankings
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Summary & Rankings", "Ranked by association with team winning (destination barthag)")

# Ranking cards
rank_data = [
    ("1", "C  —  Two-Way Engine",  ORANGE,
     "Dest. barthag 0.613  |  Barthag chg +0.104  |  WAB -4.6  |  BPM +2.55  |  n=46",
     "Rarest and most impactful. Versatile forwards dominate. Massive upward program leap."),
    ("2", "B  —  Elite Playmaker", GREEN,
     "Dest. barthag 0.593  |  Barthag chg +0.064  |  WAB -5.4  |  BPM +0.55  |  n=914",
     "Best stacking effect in the portal (2+ -> barthag 0.659). Highest-volume playable asset."),
    ("3", "A  —  Elite Rebounder", BLUE,
     "Dest. barthag 0.548  |  Barthag chg -0.005  |  WAB -6.7  |  BPM +0.13  |  n=924",
     "Still above baseline, but trade down in program quality. Stacking still helps (p=0.019)."),
]
for i, (rank, label, color, stats_line, note) in enumerate(rank_data):
    y = 1.55 + i * 1.65
    add_rect(slide, 0.35, y, 0.7, 1.45, fill_color=GOLD)
    add_text(slide, rank, 0.35, y+0.28, 0.7, 0.7,
             font_size=36, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 1.15, y, 11.53, 1.45, fill_color=color)
    add_text(slide, label,      1.3, y+0.07, 11.2, 0.48,
             font_size=20, bold=True, color=WHITE)
    add_text(slide, stats_line, 1.3, y+0.55, 11.2, 0.38,
             font_size=12, color=RGBColor(0xFF,0xEE,0xEE))
    add_text(slide, note,       1.3, y+0.95, 11.2, 0.38,
             font_size=11, italic=True, color=RGBColor(0xFF,0xDD,0xDD))

add_text(slide,
    "Surprising finding: Elite Rebounders (the largest category) move DOWN in program quality on average "
    "— suggesting bigs use the portal to seek playing time, not upward mobility.",
    0.35, 6.6, 12.63, 0.7,
    font_size=12, italic=True, color=SJU_CRIMSON)


# ===========================================================================
# SLIDE 11 — Actionable Recommendations
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, SJU_CRIMSON, "Actionable Recommendations", "For SJU Basketball — portal strategy insights")

recs = [
    (ORANGE, "Target Two-Way Engine forwards",
     "Only 46 exist in the full dataset — when identified, prioritize above all other targets. "
     "Their BPM (+2.55) and team barthag impact are class-leading."),
    (GREEN,  "Stack playmakers when possible",
     "Teams with 2+ Elite Playmakers average barthag 0.659. "
     "If you acquire one, go back to the portal for a second. The synergy is statistically significant."),
    (BLUE,   "Use Elite Rebounders as role clarity fits",
     "They move down in program quality — use that to your advantage. "
     "They fill a specific need (offensive rebounding, second-chance points) "
     "at a lower recruiting cost."),
    (DARK_GRAY, "Cross-reference position needs with category",
     "Playmakers are almost all guards; Two-Way Engines are versatile forwards. "
     "Roster construction should target the category that fills the current position gap."),
]
for i, (color, title, body) in enumerate(recs):
    y = 1.55 + i * 1.38
    add_rect(slide, 0.35, y, 0.18, 1.2, fill_color=color)
    add_rect(slide, 0.63, y, 12.05, 1.2, fill_color=LIGHT_GRAY)
    add_text(slide, title, 0.78, y+0.06, 11.7, 0.42,
             font_size=15, bold=True, color=color)
    add_text(slide, body,  0.78, y+0.5,  11.7, 0.62,
             font_size=12, color=DARK_GRAY)


# ===========================================================================
# SLIDE 12 — oreb_rate Prediction Model
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, BLUE,
            "Predicting Change in Offensive Rebound Rate",
            "Random Forest model — target: change_oreb_rate after transferring")

# Model performance cards
stat_card(slide, "MAE",         "1.619", 0.35,  1.55, w=2.3, h=1.05, bg_color=BLUE)
stat_card(slide, "R²",          "0.544", 2.85,  1.55, w=2.3, h=1.05, bg_color=BLUE)
stat_card(slide, "Algorithm",   "Random\nForest", 5.35, 1.55, w=2.5, h=1.05, bg_color=DARK_GRAY)
stat_card(slide, "Estimators",  "200",   8.05,  1.55, w=2.3, h=1.05, bg_color=DARK_GRAY)
stat_card(slide, "Train/Test",  "80/20", 10.55, 1.55, w=2.4, h=1.05, bg_color=DARK_GRAY)

# Feature importance chart
add_image(slide, f"{BASE}/chart_oreb_feature_importance.png", 0.35, 2.8, 7.2)

# Key findings box
bullet_box(slide,
    ["Previous oreb_rate is by far the #1 predictor (0.38 importance)",
     "Position matters: C and PF/C add significant signal",
     "Minutes per game (mpg) is the top non-historical feature",
     "Team offensive context (change_team_adjo) drives changes",
     "R² = 0.544 — model explains ~54% of variance in oreb change",
     "MAE of 1.62 percentage points on held-out test set"],
    7.75, 2.8, 5.2, 3.3,
    title="Key Findings", bg_color=LIGHT_GRAY, title_color=BLUE, font_size=12)

add_text(slide,
    "Takeaway: A player's prior offensive rebounding rate is the strongest signal — "
    "big men maintain their rebounding identity regardless of where they transfer. "
    "Role (minutes) and team offensive system modulate the margin of change.",
    0.35, 6.35, 12.63, 0.9,
    font_size=12, italic=True, color=DARK_GRAY)


# ===========================================================================
# SLIDE 13 — AST/TO Ratio Prediction Model
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, WHITE)
header_band(slide, GREEN,
            "Predicting Change in Assist-to-Turnover Ratio",
            "Random Forest model — target: change_ast/to after transferring")

# Model performance cards
stat_card(slide, "MAE",         "0.367", 0.35,  1.55, w=2.3, h=1.05, bg_color=GREEN)
stat_card(slide, "R²",          "0.469", 2.85,  1.55, w=2.3, h=1.05, bg_color=GREEN)
stat_card(slide, "Algorithm",   "Random\nForest", 5.35, 1.55, w=2.5, h=1.05, bg_color=DARK_GRAY)
stat_card(slide, "Estimators",  "200",   8.05,  1.55, w=2.3, h=1.05, bg_color=DARK_GRAY)
stat_card(slide, "Train/Test",  "80/20", 10.55, 1.55, w=2.4, h=1.05, bg_color=DARK_GRAY)

# Feature importance chart
add_image(slide, f"{BASE}/chart_ast_feature_importance.png", 0.35, 2.8, 7.2)

# Key findings box
bullet_box(slide,
    ["Previous ast/to is the dominant predictor (0.41 importance)",
     "Minutes per game (mpg) is the #2 feature — role determines output",
     "Destination offensive system (adj_o) influences playmaking",
     "Position tags for PG archetypes carry meaningful signal",
     "change_team_adjo: better offenses boost a guard's ast/to",
     "R² = 0.469 — harder to predict than oreb_rate (noisier skill)"],
    7.75, 2.8, 5.2, 3.3,
    title="Key Findings", bg_color=LIGHT_GRAY, title_color=GREEN, font_size=12)

add_text(slide,
    "Takeaway: Playmaking efficiency is habit-driven — prior ast/to is the clearest predictor. "
    "However, landing in a high-quality offense and earning more minutes both independently "
    "push ast/to upward, suggesting system fit is a meaningful lever for playmaker development.",
    0.35, 6.35, 12.63, 0.9,
    font_size=12, italic=True, color=DARK_GRAY)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_path = f"{BASE}/transfers_analysis.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
